# -*- coding: utf-8 -*-
"""Perky IR덱 6장 (16:9).

  python3 docs/발표/build-ir-deck.py      # 필요: pip install python-pptx

화면 캡처는 docs/이미지/{find,policies,result}.png 를 읽는다. 캡처를 새로 찍어
같은 파일명으로 덮어쓰고 이 스크립트를 다시 돌리면 덱에 반영된다.

수치(정책 상한·테스트 수·PR 수)는 손으로 적혀 있다. 레포가 바뀌면 여기도 고친다 —
발표 자료라 자동 반영보다 사람이 한 번 확인하는 편이 낫다고 봤다.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

HERE = Path(__file__).resolve().parent
OUT = str(HERE / "Perky_IR덱.pptx")
SHOTS = str(HERE.parent / "이미지")   # docs/이미지 — 화면 캡처를 바꾸면 여기만 덮어쓰면 된다
FONT = "Malgun Gothic"

C = dict(  # 앱의 design-tokens.md 팔레트를 그대로 쓴다
    brand=RGBColor(0x56,0x7c,0x8d), brandD=RGBColor(0x2c,0x45,0x52),
    accent=RGBColor(0xa8,0x54,0x34), ink=RGBColor(0x2f,0x41,0x56),
    ink6=RGBColor(0x52,0x6a,0x83), ink5=RGBColor(0x5b,0x6f,0x86),
    sand=RGBColor(0xf5,0xef,0xeb), pale=RGBColor(0xe8,0xed,0xf2),
    white=RGBColor(0xff,0xff,0xff), brand50=RGBColor(0xef,0xf4,0xf8),
)
IN = 914400
W, H = int(13.333*IN), int(7.5*IN)
M = int(0.72*IN)                 # 좌우 여백
CW = W - 2*M                     # 본문 폭

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(W), Emu(H)
BLANK = prs.slide_layouts[6]

def txt(s, x, y, w, h, lines, *, sz=14, bold=False, color=C["ink"],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=2, url=None):
    tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(spacing)
        for seg in (line if isinstance(line, list) else [line]):
            t, u = seg if isinstance(seg, tuple) else (seg, url)
            r = p.add_run(); r.text = t
            f = r.font; f.name, f.size, f.bold = FONT, Pt(sz), bold
            f.color.rgb = C["brand"] if u else color
            if u:
                r.hyperlink.address = u
                f.underline = True
    return tb

def card(s, x, y, w, h, *, fill=C["pale"], line=None, radius=0.04):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh

def head(s, kicker, title, sub=None):
    txt(s, M, int(0.52*IN), CW, int(0.3*IN), [kicker], sz=13, bold=True, color=C["accent"])
    txt(s, M, int(0.86*IN), CW, int(0.8*IN), [title], sz=36, bold=True, color=C["ink"])
    if sub:
        txt(s, M, int(1.62*IN), CW, int(0.4*IN), [sub], sz=16, color=C["ink6"])

def page(s, n):
    txt(s, W-M-int(1*IN), H-int(0.55*IN), int(1*IN), int(0.3*IN), [f"{n} / 6"],
        sz=11, color=C["ink5"], align=PP_ALIGN.RIGHT)

# ══ 1 · 문제 ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "문제", "지원금은 이미 있습니다.", "못 받는 이유는 돈이 없어서가 아닙니다.")
AG = [("국토교통부", ["청년월세 지원", "최대 480만원"]),
      ("전북특별자치도", ["전북청년 지역정착", "최대 360만원"]),
      ("익산시", ["익산형 청년월세 · 이사비", "최대 240만 / 50만원"]),
      ("국토부 · 복지부", ["청년 주거급여 분리지급", "월 최대 21만 2천원"])]
cw, gap = int((CW - 3*int(0.22*IN))/4), int(0.22*IN)
for i, (name, body) in enumerate(AG):
    x = M + i*(cw+gap)
    card(s, x, int(2.3*IN), cw, int(1.55*IN), fill=C["white"], line=C["pale"])
    txt(s, x+int(0.2*IN), int(2.5*IN), cw-int(0.4*IN), int(0.3*IN), [name], sz=15, bold=True, color=C["brandD"])
    txt(s, x+int(0.2*IN), int(2.92*IN), cw-int(0.4*IN), int(0.8*IN),
        [body[0]], sz=12, color=C["ink6"])
    txt(s, x+int(0.2*IN), int(3.35*IN), cw-int(0.4*IN), int(0.35*IN),
        [body[1]], sz=14, bold=True, color=C["accent"])
card(s, M, int(4.6*IN), CW, int(1.85*IN), fill=C["sand"])
txt(s, M+int(0.35*IN), int(4.9*IN), CW-int(0.7*IN), int(1.1*IN), [
    "네 곳이 따로 공고를 냅니다. 청년은 각각을 찾아내고, 나이·지역·소득·학적이 정책마다",
    "다르게 얽힌 조건을 자기에게 맞춰본 다음, 받아도 실제로 얼마가 남는지 또 계산해야 합니다.",
], sz=16, color=C["ink"], spacing=6)
txt(s, M+int(0.35*IN), int(5.86*IN), CW-int(0.7*IN), int(0.3*IN),
    ["→ 정작 대상인 청년이 자기가 받을 수 있는지 모릅니다."], sz=15, bold=True, color=C["accent"])
page(s, 1)

# ══ 2 · 해결 ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "해결", "두 단계로 쪼갰습니다.", "찾는 일과 계산하는 일을 섞지 않습니다.")
FLOW = [("질문 4개", ["나이 · 지역", "상태 · 소득 구간"], C["brand50"]),
        ("1층 · 발견", ["해당될 수 있는", "지원금 목록"], C["white"]),
        ("계약 조건", ["보증금 · 월세", "관리비 · 기간"], C["brand50"]),
        ("2층 · 계산", ["자격 판정 +", "실부담 주거비"], C["white"])]
bw = int(2.55*IN); ar = int(0.42*IN)
total = 4*bw + 3*ar; sx = (W-total)//2
for i, (t, body, bg) in enumerate(FLOW):
    x = sx + i*(bw+ar)
    card(s, x, int(2.35*IN), bw, int(1.5*IN), fill=bg, line=C["brand"])
    txt(s, x, int(2.6*IN), bw, int(0.35*IN), [t], sz=17, bold=True, color=C["brandD"], align=PP_ALIGN.CENTER)
    txt(s, x, int(3.06*IN), bw, int(0.7*IN), body, sz=12, color=C["ink6"], align=PP_ALIGN.CENTER, spacing=1)
    if i < 3:
        txt(s, x+bw, int(2.92*IN), ar, int(0.4*IN), ["→"], sz=20, bold=True,
            color=C["brand"], align=PP_ALIGN.CENTER)
card(s, M, int(4.7*IN), CW, int(1.72*IN), fill=C["sand"])
txt(s, M+int(0.35*IN), int(4.98*IN), CW-int(0.7*IN), int(0.35*IN),
    ["1층만 완성되어도 제품으로 성립합니다."], sz=17, bold=True, color=C["ink"])
txt(s, M+int(0.35*IN), int(5.46*IN), CW-int(0.7*IN), int(0.7*IN), [
    "계약 조건을 아직 모르는 사람도 \"내가 뭘 받을 수 있는지\"는 알고 갈 수 있어야 합니다.",
    "그래서 2층은 1층 위에 얹었고, 두 층은 따로 열립니다.",
], sz=14, color=C["ink6"], spacing=3)
page(s, 2)

# ══ 3 · 데모 ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "프로토타입", "질문 4개 · 1분.")
SHOT = [("find.png", "① 질문 4개", "답을 바꿀 때마다 건수가 바로 바뀐다"),
        ("policies.png", "② 받을 수 있는 것만", "태그와 이유, 그리고 공고 원문 링크"),
        ("result.png", "③ 실제로 낼 돈", "최대 지원 가능액과 최종 예상 주거비")]
ih = int(3.62*IN); iw = int(ih*780/1688)
gap = int(0.62*IN); total = 3*iw + 2*gap; sx = (W-total)//2
for i, (fn, cap, sub) in enumerate(SHOT):
    x = sx + i*(iw+gap)
    s.shapes.add_picture(f"{SHOTS}/{fn}", Emu(x), Emu(int(1.65*IN)), Emu(iw), Emu(ih))
    txt(s, x-int(0.2*IN), int(5.42*IN), iw+int(0.4*IN), int(0.3*IN), [cap],
        sz=15, bold=True, color=C["brandD"], align=PP_ALIGN.CENTER)
    txt(s, x-int(0.2*IN), int(5.8*IN), iw+int(0.4*IN), int(0.5*IN), [sub],
        sz=11, color=C["ink5"], align=PP_ALIGN.CENTER, spacing=1)
txt(s, M, int(6.42*IN), CW, int(0.3*IN),
    [[("wonvengers.vercel.app", "https://wonvengers.vercel.app"),
      ("  — 계정 없이 열립니다. 모바일 화면 기준으로 만들었습니다.", None)]],
    sz=13, color=C["ink6"], align=PP_ALIGN.CENTER)
page(s, 3)

# ══ 4 · 원칙 (AI 활용) ═════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "AI 활용", "판정을 AI에 맡기지 않았습니다.")
colw = (CW - int(0.4*IN))//2
LEFT = ["같은 입력은 같은 결과를 내야 합니다 (PRD F3-2).",
        "화면이 매번 모델 응답에 따라 달라지면 그 보장이 깨집니다.",
        "",
        "공고의 소득 상한·월세 상한은 산문에 섞여 있습니다.",
        "판정에 필요한 숫자를 기계가 정확히 뽑아낼 수 없습니다.",
        "",
        "추측해서 채운 숫자는 사람이 돈 계산을 틀리게 합니다."]
RIGHT = ["공고 원문을 사람이 검수해 규칙 데이터로 옮깁니다.",
         "정책 5건 전부 출처 URL과 검수 날짜를 화면에 밝힙니다.",
         "",
         "판정할 수 없는 값은 통과시키지 않고 '확인 필요' 로 남깁니다.",
         "검수 전 정책은 화면이 그렇게 말합니다.",
         "",
         "AI 는 개발과 데이터 정규화에 썼습니다 — 판정이 아니라."]
for i, (title, body, bg) in enumerate([("왜 안 썼나", LEFT, C["white"]), ("대신 무엇을 했나", RIGHT, C["brand50"])]):
    x = M + i*(colw + int(0.4*IN))
    card(s, x, int(1.72*IN), colw, int(3.5*IN), fill=bg, line=C["pale"])
    txt(s, x+int(0.32*IN), int(1.98*IN), colw-int(0.64*IN), int(0.35*IN), [title],
        sz=17, bold=True, color=C["accent"] if i == 0 else C["brandD"])
    txt(s, x+int(0.32*IN), int(2.5*IN), colw-int(0.64*IN), int(2.6*IN), body,
        sz=13, color=C["ink6"], spacing=5)
card(s, M, int(5.45*IN), CW, int(1.05*IN), fill=C["sand"])
txt(s, M, int(5.72*IN), CW, int(0.5*IN), ["“틀린 확신보다 정직한 모름이 낫다”"],
    sz=22, bold=True, color=C["accent"], align=PP_ALIGN.CENTER)
page(s, 4)

# ══ 5 · 검증 ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "검증", "숫자를 믿어도 되는지부터 확인했습니다.")
# 숫자는 2026-08-27 기준으로 다시 셌다. 테스트는 `npm test` 의 통과 건수,
# PR 은 `gh pr list --state merged`, 커밋은 `git rev-list --count origin/develop`.
STAT = [("5건", "정책 전부 공고 원문 대조"), ("261", "테스트 통과 (로직 + 화면)"),
        ("60", "PR · 커밋 115개"), ("4주", "배포까지 (계정 없이 열림)")]
cw = int((CW - 3*int(0.22*IN))/4); gap = int(0.22*IN)
for i, (n, label) in enumerate(STAT):
    x = M + i*(cw+gap)
    card(s, x, int(1.72*IN), cw, int(1.3*IN), fill=C["white"], line=C["pale"])
    txt(s, x, int(1.92*IN), cw, int(0.55*IN), [n], sz=34, bold=True, color=C["accent"], align=PP_ALIGN.CENTER)
    txt(s, x+int(0.15*IN), int(2.55*IN), cw-int(0.3*IN), int(0.4*IN), [label],
        sz=11, color=C["ink6"], align=PP_ALIGN.CENTER)
card(s, M, int(3.3*IN), CW, int(2.85*IN), fill=C["sand"])
txt(s, M+int(0.35*IN), int(3.55*IN), CW-int(0.7*IN), int(0.35*IN),
    ["대조하지 않았으면 못 잡았을 오류"], sz=13, bold=True, color=C["accent"])
txt(s, M+int(0.35*IN), int(3.95*IN), CW-int(0.7*IN), int(0.45*IN),
    ["기준 중위소득 표가 2027년 값이었습니다."], sz=22, bold=True, color=C["ink"])
txt(s, M+int(0.35*IN), int(4.55*IN), CW-int(0.7*IN), int(1.5*IN), [
    "income-thresholds.json 은 2026년 사업용인데 숫자는 2027년 기준(6.7% 인상분)이 들어 있었습니다.",
    "모든 소득 상한이 6.7% 씩 높게 계산되고 있었습니다 — 자격이 없는 사람을 통과시키는 방향의 오차입니다.",
    "보건복지부 보도자료 원문으로 1~6인 가구 값을 확인해 2026년 값으로 교체했습니다.",
    "7인 가구는 보도자료에 없어 고시 인용값을 쓰고 \"재확인 필요\"로 기록에 남겼습니다.",
], sz=14, color=C["ink6"], spacing=5)
page(s, 5)

# ══ 6 · 남은 것과 다음 ═════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
head(s, "남은 것과 다음", "아직 못 한 것을 먼저 말합니다.")
colw = (CW - int(0.4*IN))//2
card(s, M, int(1.72*IN), colw, int(2.5*IN), fill=C["white"], line=C["accent"])
txt(s, M+int(0.32*IN), int(1.96*IN), colw-int(0.64*IN), int(0.35*IN), ["남은 리스크"],
    sz=16, bold=True, color=C["accent"])
txt(s, M+int(0.32*IN), int(2.42*IN), colw-int(0.64*IN), int(1.7*IN), [
    "정책 5건 중 익산형 청년월세 1건이 팀 교차검수 전입니다.",
    "",
    "숨기지 않습니다 — 그 카드는 화면에서 \"아직 공고 원문과",
    "대조하지 않았습니다. 신청 전에 원문을 직접 확인하세요\"",
    "라고 스스로 말합니다.",
], sz=13, color=C["ink6"], spacing=4)
x2 = M + colw + int(0.4*IN)
card(s, x2, int(1.72*IN), colw, int(2.5*IN), fill=C["brand50"], line=C["pale"])
txt(s, x2+int(0.32*IN), int(1.96*IN), colw-int(0.64*IN), int(0.35*IN), ["다음"],
    sz=16, bold=True, color=C["brandD"])
txt(s, x2+int(0.32*IN), int(2.42*IN), colw-int(0.64*IN), int(1.7*IN), [
    "정책 5개 → 8~12개 (PRD 목표).",
    "",
    "온통청년 API 로 익산 주거 정책 27건을 훑어 24건을",
    "후보로 뽑아 뒀습니다. 다만 지역 필터가 새서 다른 시·도",
    "사업이 섞입니다 — 공고 원문을 열어 확인한 것만 넣습니다.",
], sz=13, color=C["ink6"], spacing=4)
card(s, M, int(4.45*IN), CW, int(1.75*IN), fill=C["sand"])
txt(s, M, int(4.72*IN), CW, int(0.5*IN), ["Perky"], sz=32, bold=True, color=C["ink"], align=PP_ALIGN.CENTER)
txt(s, M, int(5.32*IN), CW, int(0.32*IN),
    ["팀 Wonvengers  ·  원광대학교  ·  멋쟁이사자처럼 Campus AX-Ton"],
    sz=14, color=C["ink6"], align=PP_ALIGN.CENTER)
txt(s, M, int(5.72*IN), CW, int(0.32*IN),
    [[("wonvengers.vercel.app", "https://wonvengers.vercel.app"), ("      ", None),
      ("github.com/Campus-AX-Ton-wku/Wonvengers", "https://github.com/Campus-AX-Ton-wku/Wonvengers")]],
    sz=13, align=PP_ALIGN.CENTER)
page(s, 6)

prs.save(OUT)
print("saved:", OUT, f"({prs.slide_width}x{prs.slide_height})")
